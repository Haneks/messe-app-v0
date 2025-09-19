#!/usr/bin/env python3
"""
PowerPoint Image Generator with DeepAI API
Automatically generates and inserts contextually relevant images into PowerPoint slides
"""

import os
import sys
import requests
import time
import logging
from typing import List, Dict, Optional, Tuple
from pathlib import Path
import tempfile
import json

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    from PIL import Image
except ImportError as e:
    print(f"❌ Missing required libraries. Please install:")
    print("pip install python-pptx pillow requests")
    sys.exit(1)

# Configuration
class Config:
    """Configuration class for the PowerPoint Image Generator"""
    
    # DeepAI API Configuration
    DEEPAI_API_KEY = os.getenv('DEEPAI_API_KEY', 'f2c9bdf5-217a-452e-bf0f-3c5f5f5785c2')
    DEEPAI_API_URL = 'https://api.deepai.org/api/text2img'
    
    # Image Configuration
    IMAGE_WIDTH = Inches(4)  # Default image width
    IMAGE_HEIGHT = Inches(3)  # Default image height
    IMAGE_POSITION_X = Inches(6)  # X position (right side)
    IMAGE_POSITION_Y = Inches(1.5)  # Y position (top)
    
    # API Configuration
    REQUEST_TIMEOUT = 30  # seconds
    MAX_RETRIES = 3
    RETRY_DELAY = 2  # seconds between retries
    
    # Logging Configuration
    LOG_LEVEL = logging.INFO
    LOG_FORMAT = '%(asctime)s - %(levelname)s - %(message)s'

class PowerPointImageGenerator:
    """Main class for generating and inserting images into PowerPoint presentations"""
    
    def __init__(self, api_key: str = None):
        """
        Initialize the PowerPoint Image Generator
        
        Args:
            api_key (str): DeepAI API key
        """
        self.api_key = api_key or Config.DEEPAI_API_KEY
        self.setup_logging()
        self.temp_dir = tempfile.mkdtemp()
        
        if not self.api_key or self.api_key == 'your-deepai-api-key-here':
            self.logger.error("❌ DeepAI API key not configured!")
            self.logger.error("Please set DEEPAI_API_KEY environment variable or pass it to the constructor")
            raise ValueError("DeepAI API key is required")
    
    def setup_logging(self):
        """Setup logging configuration"""
        logging.basicConfig(
            level=Config.LOG_LEVEL,
            format=Config.LOG_FORMAT,
            handlers=[
                logging.StreamHandler(sys.stdout),
                logging.FileHandler('powerpoint_image_generator.log')
            ]
        )
        self.logger = logging.getLogger(__name__)
    
    def extract_slide_titles(self, presentation_path: str) -> List[Dict[str, any]]:
        """
        Extract titles from all slides in the presentation
        
        Args:
            presentation_path (str): Path to the PowerPoint file
            
        Returns:
            List[Dict]: List of slide information with titles
        """
        try:
            prs = Presentation(presentation_path)
            slides_info = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_info = {
                    'index': slide_idx,
                    'title': '',
                    'has_title': False,
                    'layout_name': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'
                }
                
                # Try to find title in different ways
                title_text = self._extract_title_from_slide(slide)
                
                if title_text:
                    slide_info['title'] = title_text
                    slide_info['has_title'] = True
                    self.logger.info(f"📄 Slide {slide_idx + 1}: '{title_text}'")
                else:
                    slide_info['title'] = f"Slide {slide_idx + 1}"
                    self.logger.warning(f"⚠️  No title found for slide {slide_idx + 1}")
                
                slides_info.append(slide_info)
            
            self.logger.info(f"✅ Extracted {len(slides_info)} slides from presentation")
            return slides_info
            
        except Exception as e:
            self.logger.error(f"❌ Error extracting slide titles: {str(e)}")
            raise
    
    def _extract_title_from_slide(self, slide) -> str:
        """
        Extract title text from a slide using multiple methods
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            str: Extracted title text
        """
        title_text = ""
        
        # Method 1: Try slide.shapes.title
        try:
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                if hasattr(slide.shapes.title, 'text'):
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        return title_text
        except:
            pass
        
        # Method 2: Look for title placeholder
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    if shape.placeholder_format.type == 1:  # Title placeholder
                        if hasattr(shape, 'text') and shape.text.strip():
                            return shape.text.strip()
        except:
            pass
        
        # Method 3: Look for the first text box with large font
        try:
            largest_font_size = 0
            largest_text = ""
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    try:
                        if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.runs:
                                    for run in paragraph.runs:
                                        if hasattr(run.font, 'size') and run.font.size:
                                            font_size = run.font.size.pt
                                            if font_size > largest_font_size:
                                                largest_font_size = font_size
                                                largest_text = shape.text.strip()
                    except:
                        continue
            
            if largest_text and largest_font_size > 20:  # Assume titles are > 20pt
                return largest_text
        except:
            pass
        
        # Method 4: First non-empty text shape
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.strip()
                    if len(text) < 100:  # Assume titles are shorter
                        return text
        except:
            pass
        
        return title_text
    
    def generate_image_prompt(self, slide_title: str, slide_index: int) -> str:
        """
        Generate an optimized prompt for image generation based on slide title
        
        Args:
            slide_title (str): Title of the slide
            slide_index (int): Index of the slide
            
        Returns:
            str: Optimized prompt for image generation
        """
        # Clean the title
        clean_title = slide_title.strip()
        
        # Detect liturgical content and enhance prompt
        liturgical_keywords = {
            'messe': 'catholic mass ceremony, church interior, altar',
            'évangile': 'gospel book, bible, religious scripture, holy light',
            'lecture': 'bible reading, scripture, religious text, church lectern',
            'psaume': 'psalm, religious music, church choir, spiritual',
            'communion': 'holy communion, eucharist, chalice, bread and wine',
            'chant': 'church choir, religious music, hymn, spiritual singing',
            'prière': 'prayer, hands in prayer, spiritual meditation, church',
            'liturgie': 'liturgical ceremony, church service, religious ritual',
            'célébration': 'religious celebration, church ceremony, festive',
            'sanctus': 'holy, sacred, church bells, divine light',
            'gloria': 'glory, heavenly light, angels, divine radiance',
            'kyrie': 'mercy, compassion, gentle light, peaceful',
            'offertoire': 'offering, gifts, altar, religious ceremony',
            'entrée': 'church entrance, procession, welcoming, gathering',
            'sortie': 'church exit, blessing, peaceful departure',
            'noël': 'christmas, nativity, star, peaceful night',
            'pâques': 'easter, resurrection, sunrise, hope, new life',
            'avent': 'advent, waiting, candles, purple, preparation'
        }
        
        # Base prompt
        base_prompt = clean_title.lower()
        enhanced_prompt = clean_title
        
        # Enhance with liturgical context
        for keyword, enhancement in liturgical_keywords.items():
            if keyword in base_prompt:
                enhanced_prompt = f"{clean_title}, {enhancement}"
                break
        
        # Add general religious/spiritual context if no specific keyword found
        if enhanced_prompt == clean_title:
            if any(word in base_prompt for word in ['dieu', 'seigneur', 'christ', 'jésus', 'marie']):
                enhanced_prompt = f"{clean_title}, religious art, spiritual, peaceful, divine light"
            else:
                enhanced_prompt = f"{clean_title}, peaceful, serene, beautiful, artistic"
        
        # Add quality modifiers
        final_prompt = f"{enhanced_prompt}, high quality, professional, clean, beautiful lighting, artistic composition"
        
        self.logger.info(f"🎨 Generated prompt for slide {slide_index + 1}: '{final_prompt}'")
        return final_prompt
    
    def generate_image_with_deepai(self, prompt: str, slide_index: int) -> Optional[str]:
        """
        Generate image using DeepAI API
        
        Args:
            prompt (str): Text prompt for image generation
            slide_index (int): Index of the slide
            
        Returns:
            Optional[str]: Path to the generated image file, or None if failed
        """
        headers = {'api-key': self.api_key}
        data = {'text': prompt}
        
        for attempt in range(Config.MAX_RETRIES):
            try:
                self.logger.info(f"🔄 Generating image for slide {slide_index + 1} (attempt {attempt + 1}/{Config.MAX_RETRIES})")
                
                response = requests.post(
                    Config.DEEPAI_API_URL,
                    headers=headers,
                    data=data,
                    timeout=Config.REQUEST_TIMEOUT
                )
                
                if response.status_code == 200:
                    result = response.json()
                    if 'output_url' in result:
                        image_url = result['output_url']
                        return self._download_image(image_url, slide_index)
                    else:
                        self.logger.error(f"❌ No output_url in API response: {result}")
                
                elif response.status_code == 429:  # Rate limit
                    self.logger.warning(f"⚠️  Rate limit hit, waiting {Config.RETRY_DELAY * 2} seconds...")
                    time.sleep(Config.RETRY_DELAY * 2)
                    continue
                
                else:
                    self.logger.error(f"❌ API request failed with status {response.status_code}: {response.text}")
                
            except requests.exceptions.Timeout:
                self.logger.warning(f"⚠️  Request timeout for slide {slide_index + 1}, retrying...")
            except requests.exceptions.RequestException as e:
                self.logger.error(f"❌ Request error for slide {slide_index + 1}: {str(e)}")
            except Exception as e:
                self.logger.error(f"❌ Unexpected error generating image for slide {slide_index + 1}: {str(e)}")
            
            if attempt < Config.MAX_RETRIES - 1:
                time.sleep(Config.RETRY_DELAY)
        
        self.logger.error(f"❌ Failed to generate image for slide {slide_index + 1} after {Config.MAX_RETRIES} attempts")
        return None
    
    def _download_image(self, image_url: str, slide_index: int) -> Optional[str]:
        """
        Download image from URL and save to temporary file
        
        Args:
            image_url (str): URL of the image to download
            slide_index (int): Index of the slide
            
        Returns:
            Optional[str]: Path to the downloaded image file
        """
        try:
            response = requests.get(image_url, timeout=Config.REQUEST_TIMEOUT)
            if response.status_code == 200:
                image_path = os.path.join(self.temp_dir, f'slide_{slide_index + 1}_image.jpg')
                
                with open(image_path, 'wb') as f:
                    f.write(response.content)
                
                # Verify image is valid
                try:
                    with Image.open(image_path) as img:
                        img.verify()
                    self.logger.info(f"✅ Downloaded image for slide {slide_index + 1}: {image_path}")
                    return image_path
                except Exception as e:
                    self.logger.error(f"❌ Downloaded image is corrupted for slide {slide_index + 1}: {str(e)}")
                    return None
            else:
                self.logger.error(f"❌ Failed to download image: HTTP {response.status_code}")
                return None
                
        except Exception as e:
            self.logger.error(f"❌ Error downloading image for slide {slide_index + 1}: {str(e)}")
            return None
    
    def insert_image_into_slide(self, presentation_path: str, slide_index: int, image_path: str) -> bool:
        """
        Insert image into specific slide
        
        Args:
            presentation_path (str): Path to the PowerPoint file
            slide_index (int): Index of the slide
            image_path (str): Path to the image file
            
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            prs = Presentation(presentation_path)
            slide = prs.slides[slide_index]
            
            # Calculate optimal position and size
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Position image on the right side, avoiding text
            img_width = min(Config.IMAGE_WIDTH, slide_width * 0.4)
            img_height = min(Config.IMAGE_HEIGHT, slide_height * 0.4)
            
            # Position: right side, vertically centered
            img_left = slide_width - img_width - Inches(0.5)
            img_top = (slide_height - img_height) / 2
            
            # Check if there's already an image in this position
            self._remove_existing_generated_images(slide)
            
            # Add the image
            picture = slide.shapes.add_picture(
                image_path,
                img_left,
                img_top,
                img_width,
                img_height
            )
            
            # Add metadata to identify our generated images
            if hasattr(picture, 'name'):
                picture.name = f"generated_image_slide_{slide_index + 1}"
            
            # Save the presentation
            prs.save(presentation_path)
            
            self.logger.info(f"✅ Inserted image into slide {slide_index + 1}")
            return True
            
        except Exception as e:
            self.logger.error(f"❌ Error inserting image into slide {slide_index + 1}: {str(e)}")
            return False
    
    def _remove_existing_generated_images(self, slide):
        """Remove any existing generated images from the slide"""
        shapes_to_remove = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'name') and shape.name and 'generated_image' in shape.name:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except:
                pass
    
    def process_presentation(self, input_path: str, output_path: str = None) -> bool:
        """
        Process entire presentation: extract titles, generate images, insert images
        
        Args:
            input_path (str): Path to input PowerPoint file
            output_path (str): Path for output file (optional)
            
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            if not os.path.exists(input_path):
                self.logger.error(f"❌ Input file not found: {input_path}")
                return False
            
            # Set output path
            if not output_path:
                base_name = os.path.splitext(input_path)[0]
                output_path = f"{base_name}_with_images.pptx"
            
            # Create a copy of the input file
            import shutil
            shutil.copy2(input_path, output_path)
            
            self.logger.info(f"🚀 Starting processing of: {input_path}")
            self.logger.info(f"📁 Output will be saved to: {output_path}")
            
            # Extract slide titles
            slides_info = self.extract_slide_titles(output_path)
            
            if not slides_info:
                self.logger.error("❌ No slides found in presentation")
                return False
            
            # Process each slide
            successful_images = 0
            total_slides = len(slides_info)
            
            for slide_info in slides_info:
                slide_index = slide_info['index']
                slide_title = slide_info['title']
                
                if not slide_info['has_title']:
                    self.logger.info(f"⏭️  Skipping slide {slide_index + 1} (no title)")
                    continue
                
                self.logger.info(f"🎨 Processing slide {slide_index + 1}/{total_slides}: '{slide_title}'")
                
                # Generate image prompt
                prompt = self.generate_image_prompt(slide_title, slide_index)
                
                # Generate image
                image_path = self.generate_image_with_deepai(prompt, slide_index)
                
                if image_path:
                    # Insert image into slide
                    if self.insert_image_into_slide(output_path, slide_index, image_path):
                        successful_images += 1
                    else:
                        self.logger.error(f"❌ Failed to insert image into slide {slide_index + 1}")
                else:
                    self.logger.error(f"❌ Failed to generate image for slide {slide_index + 1}")
                
                # Small delay to be respectful to the API
                time.sleep(1)
            
            # Summary
            self.logger.info(f"🎉 Processing complete!")
            self.logger.info(f"📊 Successfully added images to {successful_images}/{total_slides} slides")
            self.logger.info(f"💾 Enhanced presentation saved to: {output_path}")
            
            return successful_images > 0
            
        except Exception as e:
            self.logger.error(f"❌ Error processing presentation: {str(e)}")
            return False
        finally:
            # Cleanup temporary files
            self._cleanup_temp_files()
    
    def _cleanup_temp_files(self):
        """Clean up temporary files"""
        try:
            import shutil
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
                self.logger.info("🧹 Cleaned up temporary files")
        except Exception as e:
            self.logger.warning(f"⚠️  Could not clean up temp files: {str(e)}")

def main():
    """Main function for command-line usage"""
    import argparse
    
    parser = argparse.ArgumentParser(description='Generate and insert images into PowerPoint slides using DeepAI')
    parser.add_argument('input_file', help='Path to input PowerPoint file')
    parser.add_argument('-o', '--output', help='Path to output PowerPoint file')
    parser.add_argument('-k', '--api-key', help='DeepAI API key')
    parser.add_argument('-v', '--verbose', action='store_true', help='Verbose logging')
    
    args = parser.parse_args()
    
    if args.verbose:
        Config.LOG_LEVEL = logging.DEBUG
    
    try:
        # Initialize generator
        generator = PowerPointImageGenerator(api_key=args.api_key)
        
        # Process presentation
        success = generator.process_presentation(args.input_file, args.output)
        
        if success:
            print("✅ Successfully enhanced PowerPoint presentation with images!")
            sys.exit(0)
        else:
            print("❌ Failed to enhance PowerPoint presentation")
            sys.exit(1)
            
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()